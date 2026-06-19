# -*- coding: utf-8 -*-
# Copyright 2026 蜜蜂 & CoApis Contributors
# Licensed under the Apache License, Version 2.0
from __future__ import annotations
import asyncio, json, logging, os, subprocess, time
from pathlib import Path
from typing import Any
from .registry import register_tool

logger = logging.getLogger(__name__)

# Extensions supported
SUPPORTED_FORMATS = {
    ".pdf": "pdf", ".docx": "docx", ".doc": "docx",
    ".pptx": "pptx", ".ppt": "pptx",
    ".xlsx": "xlsx", ".xls": "xlsx", ".csv": "csv",
    ".txt": "text", ".md": "text", ".rst": "text",
    ".html": "text", ".htm": "text",
    ".json": "text", ".yaml": "text", ".yml": "text", ".toml": "text",
    ".py": "text", ".js": "text", ".ts": "text",
    ".java": "text", ".go": "text", ".rs": "text",
    ".xml": "text", ".sql": "text", ".sh": "text",
    ".png": "ocr", ".jpg": "ocr", ".jpeg": "ocr", ".gif": "ocr", ".webp": "ocr",
}


def _detect_format(filepath: str) -> str | None:
    ext = Path(filepath).suffix.lower()
    return SUPPORTED_FORMATS.get(ext)


async def _read_text(filepath: str, max_chars: int = 100000) -> dict[str, Any]:
    try:
        content = Path(filepath).read_text(encoding="utf-8", errors="ignore")
        truncated = len(content) > max_chars
        return {"content": content[:max_chars], "chars": len(content), "truncated": truncated, "format": "text"}
    except Exception as e:
        return {"error": str(e)}


async def _read_pdf(filepath: str) -> dict[str, Any]:
    """Try pymupdf4llm, then markitdown, then pdftotext."""
    # Try pymupdf4llm
    try:
        import pymupdf4llm
        md = pymupdf4llm.to_markdown(filepath)
        return {"content": md, "chars": len(md), "format": "pymupdf4llm"}
    except Exception:
        pass
    # Try markitdown
    try:
        proc = await asyncio.create_subprocess_exec(
            "markitdown", filepath,
            stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE
        )
        stdout, _ = await asyncio.wait_for(proc.communicate(), timeout=30)
        content = stdout.decode(errors="replace")
        if content.strip():
            return {"content": content, "chars": len(content), "format": "markitdown"}
    except Exception:
        pass
    # Try pdftotext
    try:
        proc = await asyncio.create_subprocess_exec(
            "pdftotext", filepath, "-",
            stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE
        )
        stdout, _ = await asyncio.wait_for(proc.communicate(), timeout=30)
        content = stdout.decode(errors="replace")
        if content.strip():
            return {"content": content, "chars": len(content), "format": "pdftotext"}
    except Exception:
        pass
    return {"error": "PDF 解析失败（需要 pymupdf4llm / markitdown / pdftotext）"}


async def _read_docx(filepath: str) -> dict[str, Any]:
    try:
        from docx import Document
        doc = Document(filepath)
        text = "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
        return {"content": text, "chars": len(text), "format": "python-docx"}
    except ImportError:
        pass
    try:
        proc = await asyncio.create_subprocess_exec(
            "markitdown", filepath,
            stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE
        )
        stdout, _ = await asyncio.wait_for(proc.communicate(), timeout=30)
        content = stdout.decode(errors="replace")
        if content.strip():
            return {"content": content, "chars": len(content), "format": "markitdown"}
    except Exception:
        pass
    return {"error": "DOCX 解析失败（需要 python-docx / markitdown）"}


async def _read_pptx(filepath: str) -> dict[str, Any]:
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
        content = "\n\n".join(t for t in texts if t.strip())
        return {"content": content, "chars": len(content), "format": "python-pptx"}
    except ImportError:
        pass
    return {"error": "PPTX 解析失败（需要 python-pptx）"}


async def _read_xlsx(filepath: str) -> dict[str, Any]:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
        all_text = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(str(c) if c is not None else "" for c in row)
                if row_text.strip(" |"):
                    rows.append(row_text)
            if rows:
                all_text.append(f"=== Sheet: {sheet_name} ===\n" + "\n".join(rows))
        content = "\n\n".join(all_text)
        return {"content": content, "chars": len(content), "format": "openpyxl", "sheets": len(wb.sheetnames)}
    except ImportError:
        pass
    return {"error": "XLSX 解析失败（需要 openpyxl）"}


@register_tool(
    name="doc_reader",
    description="多格式文档解析：PDF/DOCX/PPTX/XLSX/CSV/图片OCR → 结构化文本，支持自动格式检测。",
    category="builtin",
    tags=["file", "reader", "document", "ocr"],
    scene="general"
)
async def doc_reader(
    action: str = "read",
    file_path: str = "",
    content: str = "",
    max_chars: int = 100000,
    auto_detect: bool = True,
) -> dict[str, Any]:
    """多格式文档解析。

    Args:
        action: 操作类型 (read/detect/formats)
        file_path: 文件路径
        content: 直接传入文本内容（跳过文件读取）
        max_chars: 最大字符数
        auto_detect: 自动检测格式

    Returns:
        解析结果
    """
    if action == "formats":
        return {"action": "formats", "supported": SUPPORTED_FORMATS}

    if action == "detect":
        if not file_path.strip():
            return {"error": "file_path 不能为空"}
        if not Path(file_path).exists():
            return {"error": f"文件不存在: {file_path}"}
        fmt = _detect_format(file_path)
        return {"action": "detect", "file": file_path, "format": fmt, "supported": fmt is not None,
                "size_kb": round(Path(file_path).stat().st_size / 1024, 1)}

    if content.strip():
        result = {"content": content[:max_chars], "chars": len(content), "format": "text", "truncated": len(content) > max_chars}
    elif file_path.strip():
        fp = Path(file_path).resolve()
        if not fp.exists():
            return {"error": f"文件不存在: {file_path}"}
        fmt = _detect_format(str(fp)) if auto_detect else "text"
        if not fmt:
            return {"error": f"不支持的格式: {fp.suffix}"}
        if fmt == "text":
            result = await _read_text(str(fp), max_chars)
        elif fmt == "pdf":
            result = await _read_pdf(str(fp))
        elif fmt == "docx":
            result = await _read_docx(str(fp))
        elif fmt == "pptx":
            result = await _read_pptx(str(fp))
        elif fmt == "xlsx":
            result = await _read_xlsx(str(fp))
        elif fmt == "csv":
            result = await _read_text(str(fp), max_chars)
        elif fmt == "ocr":
            return {"error": "OCR 功能需要 markitdown CLI 支持"}
        else:
            result = await _read_text(str(fp), max_chars)
    else:
        return {"error": "file_path 或 content 不能都为空"}

    if "error" in result:
        return result

    return {
        "action": "read",
        "chars": result.get("chars", 0),
        "format": result.get("format", "unknown"),
        "truncated": result.get("truncated", False),
        "content": result.get("content", ""),
    }
